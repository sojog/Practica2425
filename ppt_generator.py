import yaml
import json
from pptx import Presentation
from pptx.util import Pt
import copy


def copy_slide_from_external_prs(prs, slide_no=0, template_ppt="templater.pptx"):

    # copy from external presentation all objects into the existing presentation
    external_pres = Presentation(template_ppt)

    # specify the slide you want to copy the contents from
    ext_slide = external_pres.slides[slide_no]

    # Define the layout you want to use from your generated pptx
    SLD_LAYOUT = 0
    slide_layout = prs.slide_layouts[SLD_LAYOUT]

    # create now slide, to copy contents to 
    curr_slide = prs.slides.add_slide(slide_layout)

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated

    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs


template_ppt = Presentation("templater.pptx")
output_ppt = Presentation("templater.pptx")

# titlu
print(output_ppt.slides[0].shapes[0].shape_id)
output_ppt.slides[0].shapes[0].text_frame.text  = "Test din cod"
# intrebare
print(output_ppt.slides[0].shapes[1].shape_id)
# variantele
print(output_ppt.slides[0].shapes[2].shape_id)

# copy_slide_from_external_prs(output_ppt)



with open("intrebari_reformulate.txt") as freader:
    intrebari =  freader.read().split("\n\n")

for index, intrebare in enumerate(intrebari):

    if index:
        copy_slide_from_external_prs(output_ppt)

    intrebari_lines = intrebare.splitlines()
    raspunsuri = "\n".join(intrebari_lines[1:])

    output_ppt.slides[index].shapes[0].text_frame.text  = f"Q {index + 1}"

    output_ppt.slides[index].shapes[1].text_frame.text  = f"{intrebari_lines[0]}"

    output_ppt.slides[index].shapes[2].text_frame.text  = f"{raspunsuri}"


output_ppt.save("test_final.pptx")